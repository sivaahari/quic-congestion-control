#!/usr/bin/env python3
"""make_deck_v2.py -- Phase-1 review deck, measurement-led framing.

Design rules, in response to review feedback that v1 buried the point:
  * ONE idea per slide, stated in the slide title as a plain sentence.
  * No RFC section numbers in headlines; they go in footers as provenance.
  * Every claim on a slide is either a quotation or a number we measured.
  * A viewer who reads only the titles, in order, gets the whole argument.

Output: paper/PhaseI_Review.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = "/home/sivaa/pvseed"
FIG = f"{BASE}/analysis/figures"
OUT = f"{BASE}/paper/PhaseI_Review.pptx"
os.makedirs(f"{BASE}/paper", exist_ok=True)

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6A, 0x6A, 0x6A)
BLUE = RGBColor(0x4C, 0x72, 0xB0)
RED = RGBColor(0xC4, 0x4E, 0x52)
GREEN = RGBColor(0x55, 0xA8, 0x68)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def s_():
    return prs.slides.add_slide(BLANK)


def tb(s, x, y, w, h):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def p_(tf, text, size=18, bold=False, color=INK, after=10, first=False, align=PP_ALIGN.LEFT):
    par = tf.paragraphs[0] if first else tf.add_paragraph()
    par.text = text
    par.alignment = align
    par.space_after = Pt(after)
    for r in par.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return par


def head(s, sentence, kicker=None):
    """The slide title IS the claim, written as a sentence."""
    if kicker:
        p_(tb(s, 0.75, 0.40, 12, 0.4), kicker, size=12.5, bold=True, color=BLUE, first=True)
    p_(tb(s, 0.75, 0.78, 12, 1.05), sentence, size=27, bold=True, first=True)


def foot(s, text):
    p_(tb(s, 0.75, 6.88, 12, 0.45), text, size=10.5, color=MUTED, first=True)


def pic(s, name, x, y, w):
    p = f"{FIG}/{name}"
    if os.path.exists(p):
        s.shapes.add_picture(p, Inches(x), Inches(y), width=Inches(w))
    else:
        p_(tb(s, x, y, w, 1), f"[missing: {name}]", size=12, color=RED, first=True)


# ------------------------------------------------------------------ 1 title
s = s_()
tf = tb(s, 0.9, 2.55, 11.6, 2.4)
p_(tf, "What QUIC implementations actually do", size=32, bold=True, first=True, after=2)
p_(tf, "when your phone changes network", size=32, bold=True, after=16)
p_(tf, "Phase 1 review", size=17, color=MUTED)
tf2 = tb(s, 0.9, 5.6, 11.6, 1.0)
p_(tf2, "Capstone research · Computer Networks", size=13, color=MUTED, first=True)
p_(tf2, "sivaahari · cb.sc.u4cys24055@cb.students.amrita.edu", size=13, color=MUTED)

# ------------------------------------------------------------------ 2 setup
s = s_()
head(s, "When you walk out of your house, your connection has to relearn how fast the network is.",
     kicker="THE SETTING")
tf = tb(s, 0.85, 2.25, 11.8, 4.0)
p_(tf, "QUIC keeps the connection alive when your phone switches from Wi-Fi to mobile data. "
       "That part works — it identifies you by an ID, not by your address.",
   size=19, first=True, after=20)
p_(tf, "But it still has to work out how fast to send on the new network. It cannot see the road. "
       "So it starts slow and speeds up until something breaks.", size=19, after=20)
p_(tf, "That relearning is what you feel as a stall.", size=19, bold=True, color=BLUE)

# ------------------------------------------------------------------ 3 the rule
s = s_()
head(s, "The rulebook is explicit about what should happen.", kicker="THE RULE")
box = s.shapes.add_textbox(Inches(0.95), Inches(2.35), Inches(11.4), Inches(1.6))
tfq = box.text_frame
tfq.word_wrap = True
p_(tfq, "“…an endpoint MUST immediately reset the congestion controller and round-trip time "
        "estimator for the new path to initial values.”", size=21, bold=True, color=BLUE, first=True)
tf = tb(s, 0.95, 4.2, 11.4, 2.4)
p_(tf, "In plain terms: forget how fast the old network was. Start over.", size=19, first=True, after=18)
p_(tf, "MUST — not “should”, not “may”. This is the strongest word the standard has.",
   size=19, bold=True)
foot(s, "RFC 9000 §9.4 — quoted verbatim")

# ------------------------------------------------------------------ 4 THE GAP
s = s_()
head(s, "We checked a widely-used implementation. It doesn't do this.",
     kicker="WHAT WE FOUND · 1")
pic(s, "v2_figA_spec_vs_code.png", 2.35, 1.95, 8.7)
foot(s, "picoquic — a reference QUIC implementation, 548 of its own tests passing. "
        "Verified four independent ways: the reset function exists in seven congestion-control "
        "algorithms and is called by none of them.")

# ------------------------------------------------------------------ 5 so what
s = s_()
head(s, "So we implemented the rule ourselves — and the connection stopped working.",
     kicker="WHAT WE FOUND · 2")
pic(s, "v2_figB_compliance_broke_it.png", 3.05, 1.95, 7.3)
foot(s, "Same testbed, same file, same network settings. The only change was following the rule.")

# ------------------------------------------------------------------ 6 why
s = s_()
head(s, "Because the rule has two halves, and we had only followed one.",
     kicker="WHY IT BROKE")
pic(s, "v2_figC_two_rules.png", 1.0, 1.95, 11.4)
foot(s, "This is the finding: following the rule PARTLY is worse than ignoring it entirely.")

# ------------------------------------------------------------------ 7 the question
s = s_()
head(s, "Which raises the question nobody has answered.", kicker="THE GAP")
pic(s, "v2_figD_the_question.png", 1.5, 1.95, 10.3)
foot(s, "A 2025 study measured whether servers SUPPORT migration at all. "
        "Nobody has measured what they DO to congestion control once they accept one.")

# ------------------------------------------------------------------ 8 the paper
s = s_()
head(s, "That is the paper.", kicker="WHAT WE PROPOSE")
tf = tb(s, 0.85, 2.15, 11.8, 4.4)
p_(tf, "Take five QUIC implementations and answer three questions for each:", size=19,
   bold=True, first=True, after=18)
p_(tf, "1.   Does it reset the speed estimate when the network changes, as required?",
   size=18, after=12)
p_(tf, "2.   Does it use the free measurement the protocol already takes of the new network?",
   size=18, after=12)
p_(tf, "3.   Does it correctly ignore leftover replies from the old network?",
   size=18, after=20)
p_(tf, "Then measure, in our testbed, what each choice costs.", size=19, bold=True, after=18)
p_(tf, "We already have the testbed, the measurement tools, and the first implementation's answer.",
   size=17, color=MUTED)

# ------------------------------------------------------------------ 9 bonus
s = s_()
head(s, "One more thing we noticed along the way.", kicker="A LEAD FOR LATER")
tf = tb(s, 0.85, 2.15, 11.8, 4.4)
p_(tf, "Before a connection trusts a new network, it must run a small security check — send a "
       "random number, get it echoed back.", size=19, first=True, after=18)
p_(tf, "That check is a round trip on the new network. It measures the network's speed for free.",
   size=19, after=18)
p_(tf, "In our traces it measured the new network to within 0.2% of the true value — while the "
       "connection's own estimate was three times wrong.", size=19, bold=True, color=BLUE, after=20)
p_(tf, "The standard already allows using this for one purpose. Extending it further is a natural "
       "follow-on — but it is a section of the paper, not the paper.", size=17, color=MUTED)
foot(s, "60.1 ms measured vs 60.0 ms ground truth, in the same trace where the connection believed 20.1 ms.")

# ------------------------------------------------------------------ 10 rigour
s = s_()
head(s, "We checked our instruments before trusting them.", kicker="METHOD")
pic(s, "fig3_calibration_trap.png", 0.6, 2.15, 5.9)
pic(s, "fig4_train_length.png", 6.85, 2.15, 5.9)
foot(s, "Two flaws in the emulator found and fixed before any measurement was taken. "
        "Left: a mis-set shaper lets test traffic through untouched, so you measure your own "
        "simulator. Right: short measurement bursts break without correction.")

# ------------------------------------------------------------------ 11 status
s = s_()
head(s, "Where we are.", kicker="STATUS")
tf = tb(s, 0.85, 2.05, 11.8, 4.6)
for done, txt in [
    (True, "Testbed built, calibrated, and validated"),
    (True, "First implementation surveyed — does not follow the rule"),
    (True, "Discovered that partial compliance is worse than none"),
    (True, "Measurement and analysis tooling complete"),
    (False, "Four more implementations to survey"),
    (False, "Cost of each behaviour to quantify"),
]:
    p_(tf, ("✓   " if done else "◦   ") + txt, size=19,
       bold=done, color=GREEN if done else MUTED,
       first=(txt.startswith("Testbed")), after=14)
p_(tf, "Nothing in this deck is projected. Every number is measured.", size=17,
   bold=True, color=BLUE, after=6)

# ------------------------------------------------------------------ 12 honesty
s = s_()
head(s, "What we are deliberately not claiming.", kicker="HONESTY")
tf = tb(s, 0.85, 2.15, 11.8, 4.4)
for t in [
    "We have not shown that any particular fix is worth deploying.",
    "We have not tested on real phones or real cellular networks — this is emulation.",
    "Our own compliant implementation still fails in one scenario; one cause is diagnosed, "
    "a second is not.",
    "One implementation is not a survey. Four more to go before we can generalise.",
]:
    p_(tf, "•   " + t, size=18, first=(t.startswith("We have not shown")), after=16)
p_(tf, "We would rather state the limits than have a reviewer find them.", size=18,
   bold=True, color=BLUE)

prs.save(OUT)
print(f"wrote {OUT}")
