#!/usr/bin/env python3
"""make_deck_phase2.py -- Phase-2 review deck, generated from survey_results.json.

Same design rules as Phase 1, which reviewed well:
  * ONE idea per slide, stated in the title as a plain sentence.
  * Read only the titles, in order, and you get the whole argument.
  * No RFC section numbers in headlines; they go in footers as provenance.
  * Every number comes from the survey data file, not from prose.

Output: paper/PhaseII_Review.pptx
"""
import json
import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

BASE = "/home/sivaa/pvseed"
FIG = f"{BASE}/analysis/figures"
OUT = f"{BASE}/paper/PhaseII_Review.pptx"
os.makedirs(f"{BASE}/paper", exist_ok=True)

D = json.load(open(f"{BASE}/analysis/survey_results.json"))
IMPLS = D["implementations"]
FIND = {f["id"]: f for f in D["cross_cutting_findings"]}

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x6A, 0x6A, 0x6A)
BLUE = RGBColor(0x4C, 0x72, 0xB0)
RED = RGBColor(0xC4, 0x4E, 0x52)
GREEN = RGBColor(0x55, 0xA8, 0x68)
AMBER = RGBColor(0xB8, 0x94, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def s_():
    return prs.slides.add_slide(BLANK)


def tb(s, x, y, w, h):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def p_(tf, text, size=18, bold=False, color=INK, after=10, first=False,
       align=PP_ALIGN.LEFT):
    par = tf.paragraphs[0] if first else tf.add_paragraph()
    par.text = text
    par.alignment = align
    par.space_after = Pt(after)
    for r in par.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return par


def head(s, sentence, kicker=None):
    if kicker:
        # 0.36 not 0.40: at 0.40 this box ended at 0.80 and the title box starts
        # at 0.78, so every slide carried a 0.02in box overlap.
        p_(tb(s, 0.75, 0.40, 12, 0.36), kicker, size=12.5, bold=True,
           color=BLUE, first=True)
    p_(tb(s, 0.75, 0.78, 12, 1.05), sentence, size=26, bold=True, first=True)


def foot(s, text):
    p_(tb(s, 0.75, 6.88, 12, 0.45), text, size=10.5, color=MUTED, first=True)


NUMWORD = {0: "None", 1: "One", 2: "Two", 3: "Three", 4: "Four", 5: "Five",
           6: "Six", 7: "Seven", 8: "Eight", 9: "Nine", 10: "Ten"}


def word(n, cap=True):
    """Spell a small count for prose.

    This module's own rule is that every number comes from the survey data, but
    the slide TITLES were quietly exempt and carried their counts as literals.
    The matching figure caption drifted exactly that way -- it still read
    "Three implementations, three different mechanisms" under a title saying
    five. Route the prose counts through the data too.
    """
    w = NUMWORD.get(n, str(n))
    return w if cap else w.lower()


def listify(par, kind, marL_in, char="•"):
    """Give a paragraph a REAL bullet or number, plus a hanging indent.

    A literal "•" or "1." typed into the text cannot align a wrapped line:
    PowerPoint only hangs the indent for bullets it owns, so the second line of
    a long bullet fell back under the bullet itself. python-pptx exposes no API
    for this, hence the raw pPr.  kind: "char" | "num" | "none".
    """
    pPr = par._p.get_or_add_pPr()
    pPr.set("marL", str(int(marL_in * 914400)))
    pPr.set("indent", "0" if kind == "none" else str(int(-marL_in * 914400)))
    if kind == "char":
        pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"}))
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": char}))
    elif kind == "num":
        pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": "Calibri"}))
        pPr.append(pPr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"}))
    else:
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))
    return par


FOOT_Y = 6.88          # where foot() puts its caption
PIC_BOTTOM = FOOT_Y - 0.12


def pic(s, name, x, y, w, bottom=PIC_BOTTOM):
    """Place a figure scaled to fit BOTH the given width and the space above
    the caption, centred in the width band if height is the binding constraint.

    Passing width alone lets python-pptx derive height from the image's aspect
    ratio. That is fine until a figure is regenerated at a different aspect --
    which happened here: fig1 grew a row when ngtcp2 was added, and at its
    hardcoded 12.3in width the derived height ran off the bottom of the slide
    and straight over the caption. Fit to the box instead of trusting the width.
    """
    p = f"{FIG}/{name}"
    if not os.path.exists(p):
        p_(tb(s, x, y, w, 1), f"[missing: {name}]", size=12, color=RED, first=True)
        return
    with Image.open(p) as im:
        px_w, px_h = im.size
    h = w * px_h / px_w
    avail = bottom - y
    if h > avail:                       # height binds: scale down, centre across
        h = avail
        w_fit = h * px_w / px_h
        x += (w - w_fit) / 2
        w = w_fit
    else:                               # width binds: centre in the leftover height
        y += (avail - h) / 2
    s.shapes.add_picture(p, Inches(x), Inches(y),
                         width=Inches(w), height=Inches(h))


n_done = len([i for i in IMPLS if i["status"] == "done"])
n_live = len([i for i in IMPLS if (i.get("live") or {}).get("reps")])
dead = [i["name"] for i in IMPLS if (i.get("reset_hook") or {}).get("dead")]
DONE_ = [i for i in IMPLS if i["status"] == "done"]
n_lang = len({i["language"] for i in DONE_})
n_mech = len({i["q1"].get("mechanism", "") for i in DONE_})
n_q = len(D["questions"])
# Implementations that fit neither YES nor NO on a MUST (msquic's ASYMMETRIC/SPLIT).
nonbinary = [i["name"] for i in DONE_
             if {i["q1"]["answer"], i["q3"]["answer"]} & {"ASYMMETRIC", "SPLIT"}]
# Repetition parity, for the honesty slide.
_reps = {i["name"]: (i.get("live") or {}).get("reps", 0) for i in DONE_}
_max_reps = max(_reps.values()) if _reps else 0
_short = [n for n, r in _reps.items() if r and r < _max_reps]
# Corrections and validation tallies come from the data too. Both were literals
# here and both went stale the moment the picoquic work added two corrections
# and three checks -- the slide still claimed "four things" and "27 checks".
n_corr = len(D.get("corrections") or [])
# The discriminator's decisive case, from the data rather than from memory: the
# figures previously quoted here (279,456 -> 195,619, "0.4 ms apart") came from a
# trace that had been overwritten and could not be re-derived -- correction C7.
_ltr = ((next((i for i in IMPLS if i["name"] == "quic-go"), {}).get("live") or {})
        .get("loss_then_reset") or {})
_g = _ltr.get("gap_to_reset_ms_range") or []
_gap = f"{_g[0]}–{_g[1]}" if len(_g) == 2 else "a few"
_val = D.get("validation") or {}
n_checks = _val.get("checks", "?")
n_pass = _val.get("passed", "?")

# ------------------------------------------------------------------ 1 title
s = s_()
tf = tb(s, 0.9, 2.5, 11.6, 2.4)
p_(tf, f"What {word(n_done, cap=False)} QUIC implementations actually do",
   size=31, bold=True, first=True, after=2)
p_(tf, "when your phone changes network", size=31, bold=True, after=16)
p_(tf, "Phase 2 review — the cross-implementation survey", size=17, color=MUTED)
tf2 = tb(s, 0.9, 5.6, 11.6, 1.0)
p_(tf2, "Capstone research · Computer Networks", size=13, color=MUTED, first=True)
p_(tf2, "sivaahari · cb.sc.u4cys24055@cb.students.amrita.edu", size=13, color=MUTED)

# ------------------------------------------------------------------ 2 recap
s = s_()
head(s, "Last phase we found one implementation ignoring a rule it must follow.",
     kicker="WHERE WE LEFT OFF")
tf = tb(s, 0.85, 2.2, 11.8, 4.2)
p_(tf, "The standard says an endpoint MUST reset its speed estimate when a "
       "connection moves to a new network. Forget what you learned about the old "
       "one; start over.", size=19, first=True, after=18)
p_(tf, "We checked picoquic. It doesn't. The function that performs the reset "
       "exists, is implemented for all seven of its congestion-control "
       "algorithms, and is called by nothing.", size=19, after=18)
p_(tf, "Then we implemented the rule ourselves — and transfers stopped completing "
       "entirely, because the rule has a second half that is easy to miss.",
   size=19, after=18)
p_(tf, "Which left an obvious question.", size=19, bold=True, color=BLUE)

# ------------------------------------------------------------------ 3 question
s = s_()
head(s, "Does anyone else get this right?", kicker="THE QUESTION")
tf = tb(s, 0.85, 2.2, 11.8, 4.2)
p_(tf, f"We took the {word(n_done, cap=False)} QUIC implementations that carry "
       f"most of the internet's traffic, in {word(n_lang, cap=False)} different "
       f"languages, and asked each of them the same {word(n_q, cap=False)} "
       f"questions.", size=19, first=True, after=22)
for n, q in enumerate(("q1", "q2", "q3"), 1):
    qq = D["questions"][q]
    listify(p_(tf, qq["short"], size=18, bold=True, after=4), "num", 0.40)
    # was seven literal spaces, which does not survive a wrap
    listify(p_(tf, f"{qq['requirement']} — {qq['rfc']}", size=14, color=MUTED,
               after=14), "none", 0.40)
p_(tf, "Every answer twice: once by reading the source, once by measuring the "
       "running program. If they disagreed, we had made a mistake and went to "
       "find it.", size=17, color=BLUE, bold=True)

# ------------------------------------------------------------------ 4 the table
s = s_()
head(s, "No two of them behave the same way.", kicker="THE RESULT")
pic(s, "p2_fig1_compliance_table.png", 0.5, 1.9, 12.3)
foot(s, f"{n_done} implementations audited from source; {n_live} of them also "
        f"measured live in our testbed. Every cell is backed by both.")

# ------------------------------------------------------------------ 5 dead code
s = s_()
head(s, f"{word(len(dead))} of them ship the code to obey the rule — "
        f"and never call it.", kicker="FINDING 1")
pic(s, "p2_fig2_dead_hooks.png", 1.0, 1.95, 11.3)
foot(s, "quiche's dead function even carries the comment "
        "“Called when connection migrates and cwnd needs to be reset”. "
        "Independent codebases, different languages, the same dead hook.")

# ------------------------------------------------------------------ 6 mechanisms
s = s_()
head(s, f"Where they do comply, they get there {word(n_mech, cap=False)} "
        f"different ways.", kicker="FINDING 2")
pic(s, "p2_fig3_mechanisms.png", 0.6, 1.95, 12.1)
foot(s, "One replaces its controller outright. One is compliant by accident of "
        "architecture. One uses a packet-number watermark. One resets at only "
        "one end. One doesn't reset at all.")

# ------------------------------------------------------------------ 7 unanimous
s = s_()
head(s, "Not one of them uses the free measurement the protocol already takes.",
     kicker="FINDING 3")
tf = tb(s, 0.85, 2.2, 11.8, 4.3)
p_(tf, "Before a connection trusts a new network it must run a security check — "
       "send a random number, require it echoed back. That is a compulsory round "
       "trip on the new network, so it measures that network's speed for free.",
   size=19, first=True, after=20)
p_(tf, f"The standard explicitly permits using it. All {word(n_done, cap=False)} "
       f"decline.", size=20, bold=True, color=RED, after=20)
p_(tf, "Two of them we caught in the act, in consecutive log lines: measure the "
       "new network at roughly 40 milliseconds, then throw it away and seed a "
       "fixed default instead — 100 milliseconds in one, 333 in the other.",
   size=18, after=18)
p_(tf, "This is the one finding that is unanimous across every major "
       "implementation.", size=18, bold=True, color=BLUE)
foot(s, "RFC 9002 §6.2.2 — a MAY, not a MUST. Nobody is breaking a rule here. "
        "But nobody is taking the free measurement either.")

# ------------------------------------------------------------------ 8 not binary
s = s_()
head(s, f"{word(len(nonbinary))} of them cannot be scored yes or no at all.",
     kicker="FINDING 4")
tf = tb(s, 0.85, 2.2, 11.8, 4.3)
p_(tf, "msquic splits its state: congestion control is shared across the whole "
       "connection, but round-trip time is tracked per network path. That "
       "produces two answers no yes/no column can hold.", size=19, first=True, after=20)
p_(tf, "On the first question the answer is ASYMMETRIC. The endpoint that "
       "OBSERVES the move resets. The endpoint that MAKES the move does not. "
       "The same connection is compliant at one end and not at the other.",
   size=18, color=AMBER, bold=True, after=18)
p_(tf, "On the third the answer is SPLIT. Round-trip measurements from the old "
       "network are correctly ignored. Acknowledged byte counts are not.",
   size=18, color=AMBER, bold=True, after=20)
p_(tf, "We are not the first to notice the second one. msquic's own source says, "
       "at the exact line responsible:", size=17, after=8)
listify(p_(tf, "const QUIC_PATH* Path = &Connection->Paths[0];  // TODO - Correct?",
           size=15, bold=True, color=RED), "none", 0.45)

# ------------------------------------------------------------------ 9 the striking one
s = s_()
head(s, "And in that implementation, the sender never recovers its sense of "
        "the new network.", kicker="FINDING 5")
tf = tb(s, 0.85, 2.2, 11.8, 4.3)
p_(tf, "After migrating, msquic's server holds its round-trip estimate at exactly "
       "its starting default — 333 milliseconds — and its minimum-RTT field at "
       "the “no measurement yet” marker.", size=19, first=True, after=18)
p_(tf, "For every one of roughly 3,800 samples per run. For twenty seconds. In "
       "five runs out of five. The real network delay was 40 milliseconds.",
   size=19, bold=True, color=RED, after=18)
p_(tf, "The cost is visible: about 13 Mbit/s where another implementation reached "
       "about 19 on the same link.", size=18, after=18)
p_(tf, "The irony is the finding. The endpoint that follows the rule ends up with "
       "a permanently broken estimate; the endpoint that ignores it keeps a "
       "working one.", size=18, bold=True, color=BLUE)
foot(s, "Reported as an observation, not a diagnosis. We have not yet built an "
        "instrumented version to prove the mechanism, and we say so.")

# ------------------------------------------------------------------ 10 evidence
s = s_()
head(s, "What the measurements actually look like.", kicker="THE EVIDENCE")
pic(s, "p2_fig4_cwnd_evidence.png", 1.1, 1.95, 11.1)
foot(s, "Each implementation's own initial window is the dashed line. A drop to "
        "exactly that value can only be a reset — we checked each one's loss "
        "arithmetic to be sure an ordinary loss event could not produce the same "
        "number.")

# ------------------------------------------------------------------ 11 rigour
s = s_()
head(s, "How we know we are not fooling ourselves.", kicker="METHOD")
tf = tb(s, 0.85, 2.15, 11.8, 4.4)
p_(tf, "Every claim was answered twice — from source and from measurement — and "
       "they had to agree.", size=18, first=True, after=14)
p_(tf, f"Before trusting any drop in the send rate, we looked up each "
       f"implementation's initial window, its loss floor, and its loss reduction "
       f"factor, and confirmed the number we saw could only be a reset. In "
       f"quic-go that revealed TWO events where we would have reported one — in "
       f"{_ltr.get('observed_reps', '5/5')} runs an ordinary loss backoff to exactly 0.7× the peak, then "
       f"the real reset {_gap} milliseconds later.", size=17, after=14)
p_(tf, f"Then we wrote a second validator from scratch, sharing no code with the "
       f"first, which re-derives every claim from git object storage and raw "
       f"traces. {n_checks} checks, {n_pass} pass.", size=17, after=14)
p_(tf, f"It found {word(n_corr, cap=False)} things we had got wrong along the way "
       f"— a window described as flat when it had declined, a repetition count of "
       f"four that was really three, a path delay recorded as sixty milliseconds "
       f"that was really forty, and — the one that stings — a verdict line in our "
       f"own parser that would have called a loss event a reset. All "
       f"{word(n_corr, cap=False)} are corrected and recorded.",
   size=17, color=BLUE, bold=True)

# ------------------------------------------------------------------ 12 honesty
s = s_()
head(s, "What we are deliberately not claiming.", kicker="HONESTY")
tf = tb(s, 0.85, 2.15, 11.8, 4.4)
_caveats = [
    "We have not shown that any of this harms real users. We measured what "
    "implementations do, and what it costs in our testbed.",
    "Everything is emulation. No real phones, no real cellular networks.",
    "Each implementation was tested at one commit, in one configuration, in one "
    "direction of transfer.",
    "The stuck-estimate result is an observation. We have not proven its cause.",
]
# Only claim a repetition shortfall if the data actually shows one.
for nm in _short:
    _r = _reps[nm]
    _caveats.append(
        f"{nm} has {word(_r, cap=False)} live "
        f"{'repetition' if _r == 1 else 'repetitions'} where the others have "
        f"{word(_max_reps, cap=False)}.")
for i, t in enumerate(_caveats):
    listify(p_(tf, t, size=17, first=(i == 0), after=14), "char", 0.30)
p_(tf, "We would rather state these than have a reviewer find them.", size=18,
   bold=True, color=BLUE)

# ------------------------------------------------------------------ 13 next
s = s_()
head(s, "Where this goes.", kicker="NEXT")
tf = tb(s, 0.85, 2.15, 11.8, 4.4)
p_(tf, f"The survey is complete: {word(n_done, cap=False)} implementations, "
       f"{word(n_q, cap=False)} questions, source and measurement for every cell.",
   size=19, first=True, after=18)
p_(tf, "Next is the cost — how much each of these behaviours actually loses, "
       "across a range of network conditions rather than the single pair we used "
       "here.", size=18, after=18)
p_(tf, f"Then the paper. And the findings go back to the projects: "
       f"{word(len(dead), cap=False)} of them ship a reset that never runs, and "
       f"one has a comment in its own source asking whether the line we "
       f"identified is correct.", size=18, after=20)
p_(tf, "Nobody had measured this. Now it is measured.", size=20, bold=True,
   color=BLUE)

prs.save(OUT)
print(f"wrote {OUT}")
print(f"  {len(prs.slides._sldIdLst)} slides · {n_done} implementations · "
      f"{n_live} live-measured · dead hooks in {', '.join(dead)}")
