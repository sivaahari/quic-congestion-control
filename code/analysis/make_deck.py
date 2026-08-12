#!/usr/bin/env python3
"""make_deck.py -- generate the Phase-1 review deck (.pptx).

Every number in this deck is traceable to captured data under results/ or to a
verbatim RFC quotation. Nothing is illustrative or projected.

Output: paper/PhaseI_PVSeed_Review.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = "/home/sivaa/pvseed"
FIG = f"{BASE}/analysis/figures"
OUT = f"{BASE}/paper/PhaseI_PVSeed_Review.pptx"
os.makedirs(f"{BASE}/paper", exist_ok=True)

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
ACCENT = RGBColor(0x4C, 0x72, 0xB0)
ALERT = RGBColor(0xC4, 0x4E, 0x52)
GOOD = RGBColor(0x55, 0xA8, 0x68)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def tb(s, x, y, w, h):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size=18, bold=False, color=INK, space_after=8, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return p


def title_slide(s, title, subtitle=None, kicker=None):
    if kicker:
        t = tb(s, 0.7, 0.42, 12, 0.5)
        para(t, kicker, size=13, bold=True, color=ACCENT, first=True)
    t = tb(s, 0.7, 0.82, 12, 1.0)
    para(t, title, size=30, bold=True, first=True)
    if subtitle:
        t2 = tb(s, 0.7, 1.72, 12, 0.6)
        para(t2, subtitle, size=15, color=MUTED, first=True)


def pic(s, name, x, y, w):
    p = f"{FIG}/{name}"
    if os.path.exists(p):
        s.shapes.add_picture(p, Inches(x), Inches(y), width=Inches(w))
    else:
        tf = tb(s, x, y, w, 1)
        para(tf, f"[missing figure: {name}]", size=12, color=ALERT, first=True)


def footer(s, text):
    tf = tb(s, 0.7, 6.92, 12, 0.4)
    para(tf, text, size=10.5, color=MUTED, first=True)


# ---------------------------------------------------------------- 1. title
s = slide()
tf = tb(s, 0.9, 2.5, 11.5, 2.2)
para(tf, "Turning QUIC's Path Validation into a", size=30, bold=True, first=True)
para(tf, "Congestion-Control Measurement Primitive", size=30, bold=True)
para(tf, "Phase 1 review — establishing that the gap is real", size=17, color=MUTED, space_after=4)
tf2 = tb(s, 0.9, 5.3, 11.5, 1.2)
para(tf2, "Capstone research project · Computer Networks", size=13, color=MUTED, first=True)
para(tf2, "sivaahari · cb.sc.u4cys24055@cb.students.amrita.edu", size=13, color=MUTED)

# ---------------------------------------------------------------- 2. the idea
s = slide()
title_slide(s, "The thirty-second version", kicker="ORIENTATION")
tf = tb(s, 0.8, 2.0, 11.8, 4.4)
para(tf, "When your phone switches from Wi-Fi to mobile data, the connection survives — "
         "QUIC identifies you by a Connection ID, not an address.", size=18, first=True, space_after=16)
para(tf, "But the connection still has to relearn how fast the new network is. It starts from "
         "a crawl and ramps up. You feel that as a stall.", size=18, space_after=16)
para(tf, "It turns out the protocol already runs a compulsory speed test on the new network a "
         "moment earlier — for security reasons — and then the rulebook requires it to throw the "
         "result away.", size=18, bold=True, color=ACCENT, space_after=16)
para(tf, "Phase 1 asks one question: is that real?", size=18, space_after=10)
para(tf, "Answer: yes — on paper, in deployed code, and in measurement.", size=18, bold=True, color=GOOD)

# ---------------------------------------------------------------- 3. mechanism
s = slide()
title_slide(s, "How QUIC migration works", "The badge, and the sealed envelope", kicker="BLOCK A · THE GAP ON PAPER")
tf = tb(s, 0.8, 2.1, 11.8, 4.2)
para(tf, "1.  The client changes address. It keeps its Connection ID — the badge.", size=18, first=True, space_after=14)
para(tf, "2.  The server cannot simply trust the badge; anyone could steal one. So it sends a "
         "PATH_CHALLENGE containing random bytes, and the client must echo them back in a "
         "PATH_RESPONSE.", size=18, space_after=14)
para(tf, "3.  Only then does the server accept the new address.", size=18, space_after=20)
para(tf, "That exchange is compulsory, happens at every handover, and is a full round trip "
         "down the exact path about to be used.", size=18, bold=True, space_after=10)
para(tf, "RFC 9000 §8.2 additionally requires BOTH the challenge and the response datagrams to be "
         "padded to at least 1200 bytes. So it is a full-size packet, making a full round trip, "
         "on the new path — before any application data is risked.", size=16, color=ACCENT)
footer(s, "RFC 9000 §8.2.1, §8.2.2 — verified verbatim against the RFC text")

# ---------------------------------------------------------------- 4. the MUST
s = slide()
title_slide(s, "…and the rulebook requires the result to be discarded", kicker="BLOCK A · THE GAP ON PAPER")
box = s.shapes.add_textbox(Inches(0.9), Inches(2.05), Inches(11.6), Inches(1.9))
tf = box.text_frame
tf.word_wrap = True
para(tf, "“On confirming a peer's ownership of its new address, an endpoint MUST immediately "
         "reset the congestion controller and round-trip time estimator for the new path to "
         "initial values … unless the only change in the peer's address is its port number.”",
     size=18, bold=True, color=ALERT, first=True)
tf2 = tb(s, 0.9, 4.1, 11.6, 2.6)
para(tf2, "— RFC 9000 §9.4", size=14, color=MUTED, first=True, space_after=18)
para(tf2, "It is a MUST, not a SHOULD. The amnesia is specified, not an oversight.", size=18, space_after=12)
para(tf2, "The normative object is the phrase “initial values”. The RFC fixes THAT the controller "
          "is reset — it does not fix what those initial values must numerically be for a path "
          "about which evidence exists.", size=18, color=ACCENT)
footer(s, "This is the seam the entire project works in.")

# ---------------------------------------------------------------- 5. money slide A
s = slide()
title_slide(s, "The state of the art refuses to help here",
            "Careful Resume — IETF Proposed Standard for jump-starting congestion control",
            kicker="BLOCK A · THE GAP ON PAPER")
box = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.6), Inches(1.4))
tf = box.text_frame
tf.word_wrap = True
para(tf, "“…or the sender receives a signal from the local stack indicating that the path is now "
         "different to the observed path, the sender MUST stop using Careful Resume.”",
     size=18, bold=True, color=ALERT, first=True)
tf2 = tb(s, 0.9, 4.0, 11.6, 2.7)
para(tf2, "— draft-ietf-tsvwg-careful-resume", size=14, color=MUTED, first=True, space_after=16)
para(tf2, "Careful Resume switches itself OFF the moment the path changes — which is reasonable, "
          "because it seeds from REMEMBERED state, and on a new path an old memory is a guess, "
          "not knowledge.", size=17, space_after=12)
para(tf2, "It seeds across TIME (previous connection → new connection, same path).", size=17, space_after=6)
para(tf2, "Nothing seeds across SPACE (same connection, new path).", size=17, bold=True, color=ACCENT, space_after=12)
para(tf2, "Our answer is not to remember. It is to MEASURE — using evidence the protocol is "
          "already compelled to generate.", size=17, bold=True)

# ---------------------------------------------------------------- 6. block B
s = slide()
title_slide(s, "Does deployed code actually do this?",
            "picoquic — a widely used reference implementation, 548 tests passing",
            kicker="BLOCK B · THE GAP IN CODE")
tf = tb(s, 0.8, 2.05, 11.8, 4.5)
para(tf, "Finding: picoquic does NOT implement the §9.4 MUST-reset. Verified four independent ways:",
     size=18, bold=True, first=True, space_after=16)
para(tf, "1.  The reset hook is dead. picoquic_congestion_notification_reset has handlers in "
         "newreno, cubic, bbr, bbr1, c4, fastcc and prague — and ZERO callers.", size=16, space_after=11)
para(tf, "2.  On path validation it updates RTT, promotes the tuple, and resets the MTU. "
         "Congestion control is never touched.", size=16, space_after=11)
para(tf, "3.  cwnd is re-initialised only in CC-init and picoquic_create_path — neither on the "
         "migration route. Migration creates a new tuple on the EXISTING path object, so the "
         "window survives by construction.", size=16, space_after=11)
para(tf, "4.  Live qlog across a real IP-changing migration: cwnd flat at 241,174 bytes, never "
         "near the 15,360-byte initial window.", size=16, space_after=16)
para(tf, "So deployed reality is blind CARRY-OVER — precisely what RFC 9000 §9.4 warns "
         "“could transmit too aggressively”.", size=17, bold=True, color=ALERT)
footer(s, "picoquic 0dc8ba8b · static analysis and live measurement agree")

# ---------------------------------------------------------------- 7. fig1
s = slide()
title_slide(s, "Both behaviours, measured", kicker="BLOCK B · THE GAP IN CODE")
pic(s, "fig1_reset_vs_carryover.png", 1.5, 1.95, 10.3)
footer(s, "We implemented the compliant reset ourselves (runtime-gated, default off; 548/548 tests "
          "still pass with it disabled) — because the stack does not provide one.")

# ---------------------------------------------------------------- 8. methodology
s = slide()
title_slide(s, "Validating the instrument before trusting it",
            "Two artefacts found that would have silently invalidated every downstream number",
            kicker="METHODOLOGY")
pic(s, "fig3_calibration_trap.png", 0.55, 2.15, 6.1)
pic(s, "fig4_train_length.png", 6.85, 2.15, 6.1)
footer(s, "Left: a too-large token bucket lets a probe train through unshaped — you then measure "
          "the emulator, not the link. Right: short probe trains fail without leading-gap correction. "
          "Operating point: burst 3200 B, K ≥ 8, discard 2.")

# ---------------------------------------------------------------- 9. THE money slide
s = slide()
title_slide(s, "The measurement was there all along",
            "Same trace. Same instant. Two numbers.", kicker="BLOCK C/D · THE PAYOFF")
pic(s, "fig2_discarded_measurement.png", 0.45, 1.95, 12.4)
footer(s, "Path validation measured the new path to within 0.2% of ground truth — and the stack "
          "discarded it, while its congestion controller believed a 3× wrong value carried over "
          "from the old path.")

# ---------------------------------------------------------------- 10. numbers
s = slide()
title_slide(s, "The headline numbers", kicker="BLOCK C/D · THE PAYOFF")
rows = [
    ("What path validation measured on the new path", "60.1 ms", GOOD),
    ("Ground truth (configured)", "60.0 ms", GOOD),
    ("Error of the discarded measurement", "0.2 %", GOOD),
    ("What the congestion controller believed instead", "20.1 ms", ALERT),
    ("Peak error of the estimator while relearning", "1696 ms  (28× truth)", ALERT),
]
y = 2.35
for label, val, col in rows:
    tf = tb(s, 0.9, y, 8.0, 0.55)
    para(tf, label, size=17, first=True)
    tf2 = tb(s, 9.1, y, 3.6, 0.55)
    para(tf2, val, size=19, bold=True, color=col, first=True)
    y += 0.72
tf = tb(s, 0.9, y + 0.35, 11.8, 1.2)
para(tf, "The clean measurement is not hypothetical, not expensive, and not new work — "
         "the protocol is already required to perform it.", size=18, bold=True, color=ACCENT, first=True)

# ---------------------------------------------------------------- 11. honesty
s = slide()
title_slide(s, "What went wrong, and what it taught us", kicker="HONEST REPORTING")
tf = tb(s, 0.8, 2.05, 11.8, 4.6)
para(tf, "Our first spec-compliant reset made things WORSE — transfers stopped completing entirely.",
     size=18, bold=True, color=ALERT, first=True, space_after=14)
para(tf, "Cause: we implemented §9.4's second paragraph (reset to initial values) and missed its "
         "first — “Packets sent on the old path MUST NOT contribute to congestion control or RTT "
         "estimation for the new path.”", size=16, space_after=12)
para(tf, "Clearing the estimator let the very next in-flight OLD-path ACK re-seed it: 1.8 ms after "
         "resetting, smoothed_rtt read 20.3 ms on a path whose true RTT is 60 ms. Slow start then "
         "paced 3× too fast and destroyed the connection.", size=16, space_after=16)
para(tf, "Two things follow:", size=17, bold=True, space_after=10)
para(tf, "•  Partial §9.4 compliance is worse than none — a finding worth reporting in its own right.",
     size=16, space_after=8)
para(tf, "•  It sharpens the motivation: after migration, EVERY ordinary RTT sample is polluted by "
         "in-flight old-path packets. The path-validation exchange is, by construction, the only "
         "clean one — sent on the new path, answered on the new path.", size=16, bold=True, color=ACCENT)

# ---------------------------------------------------------------- 12. what we build
s = slide()
title_slide(s, "What Phase 2 builds", "PV-Seed — five steps", kicker="THE PLAN")
tf = tb(s, 0.8, 2.05, 11.8, 4.5)
for n, (h, d) in enumerate([
    ("PROBE", "send a short train of the padded challenges the protocol already sends"),
    ("MEASURE", "round-trip time gives the path's length; spacing between replies gives its width"),
    ("CHECK", "do the estimates agree? if they scatter, discard them and fall back"),
    ("SEED", "initialise the congestion controller from the measurement instead of from nothing"),
    ("WATCH", "send paced, and retreat to the specification's behaviour at the first sign of trouble"),
], 1):
    para(tf, f"{n}.  {h}", size=17, bold=True, color=ACCENT, first=(n == 1), space_after=2)
    para(tf, f"      {d}", size=16, space_after=13)
para(tf, "Follow the fallbacks down and you land exactly on what the RFC already demands. "
         "The worst case of PV-Seed is the status quo.", size=17, bold=True, color=GOOD)
footer(s, "picoquic already contains a live congestion-window seeding pathway (used by Careful Resume). "
          "We re-source it from measurement rather than memory — we are not adding new plumbing.")

# ---------------------------------------------------------------- 13. risk
s = slide()
title_slide(s, "What could still sink this, and how we would know early", kicker="THE PLAN")
tf = tb(s, 0.8, 2.05, 11.8, 4.5)
for r, m in [
    ("The reflected reply train may only reveal the NARROWER direction of the path",
     "proved analytically; drives a one-frame extension for the download case"),
    ("Receiver batching may destroy the reply spacing we measure",
     "dedicated experiment; a negative result here is publishable and motivates the extension"),
    ("Anti-amplification limits how many probes a server may send",
     "K ≥ 8 needs client-side priming — Phase 0 already made this mandatory, not optional"),
    ("The benefit may be small for short transfers or low-bandwidth paths",
     "we will state the regime where it matters, and the regime where it does not"),
]:
    para(tf, f"•  {r}", size=16.5, bold=True, first=(r.startswith("The reflected")), space_after=3)
    para(tf, f"     → {m}", size=15.5, color=MUTED, space_after=13)
para(tf, "Kill criterion: if the estimator cannot reach usable accuracy on clean emulated links, "
         "we pivot to a measurement paper and keep every hour invested.", size=16, bold=True)

# ---------------------------------------------------------------- 14. close
s = slide()
title_slide(s, "Phase 1 verdict", kicker="SUMMARY")
tf = tb(s, 0.8, 2.15, 11.8, 4.3)
for k, v, c in [
    ("The gap exists on paper", "RFC 9000 §9.4 mandates the discard; Careful Resume explicitly refuses to help at a path change", GOOD),
    ("The gap exists in code", "picoquic does not implement the MUST — verified four ways", GOOD),
    ("The discarded measurement is accurate", "60.1 ms vs 60.0 ms ground truth — 0.2% error", GOOD),
    ("The controller's alternative is bad", "believed 20.1 ms; peaked at 28× the truth while relearning", GOOD),
]:
    para(tf, f"✓  {k}", size=18, bold=True, color=c, first=(k == "The gap exists on paper"), space_after=3)
    para(tf, f"     {v}", size=15.5, color=MUTED, space_after=15)
para(tf, "Phase 2 begins: build the estimator, seed the controller, keep it safe.",
     size=19, bold=True, color=ACCENT)

prs.save(OUT)
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
