#!/usr/bin/env python3
"""Geometric QA on a .pptx: measure the defects that are actually user-visible.

A render shows what it looks like; this measures it. Checks:
  - anything past the slide edge, or inside the 0.4" margin
  - overlapping shapes (text drawn through other content)
  - text too long for its box, estimated from font size and box width
  - images stretched away from their native aspect ratio
"""
import sys
from pptx import Presentation

EMU_IN = 914400.0
MARGIN_IN = 0.4
FOOTER_MARGIN_IN = 0.15   # small muted text in the bottom band, by design
PICTURE = 13  # MSO_SHAPE_TYPE.PICTURE


def inches(v):
    return (v or 0) / EMU_IN


def shape_box(sh):
    try:
        return (inches(sh.left), inches(sh.top),
                inches(sh.width), inches(sh.height))
    except Exception:
        return None


def overlap(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ox = min(ax + aw, bx + bw) - max(ax, bx)
    oy = min(ay + ah, by + bh) - max(ay, by)
    return (ox, oy) if ox > 0.02 and oy > 0.02 else None


def est_text_height(shape, box_w_in):
    """Rough wrapped-height estimate. Calibri averages ~0.48 em per character."""
    total = 0.0
    for para in shape.text_frame.paragraphs:
        runs = para.runs
        if not runs:
            total += 0.18
            continue
        size_pt = next((r.font.size.pt for r in runs if r.font.size is not None), None)
        if size_pt is None:
            size_pt = para.font.size.pt if para.font.size else 18.0
        text = "".join(r.text for r in runs)
        char_w_in = (size_pt * 0.48) / 72.0
        per_line = max(int(max(box_w_in - 0.2, 0.5) / char_w_in), 1)
        lines = max(1, -(-len(text) // per_line))
        total += lines * (size_pt * 1.22) / 72.0
    return total


def main(path):
    prs = Presentation(path)
    sw, sh_in = inches(prs.slide_width), inches(prs.slide_height)
    ratio = sw / sh_in
    shape_name = ("16:9" if abs(ratio - 16 / 9) < 0.02
                  else "4:3" if abs(ratio - 4 / 3) < 0.02 else f"custom {ratio:.3f}")
    print(f"deck        : {path}")
    print(f"slide size  : {sw:.3f} x {sh_in:.3f} in  ({shape_name})")
    print(f"slides      : {len(prs.slides)}")
    print()

    issues = 0
    for idx, slide in enumerate(prs.slides, 1):
        boxes, problems = [], []
        for shp in slide.shapes:
            box = shape_box(shp)
            if box is None:
                continue
            x, y, w, h = box
            has_text = shp.has_text_frame and shp.text_frame.text.strip()
            label = (shp.text_frame.text[:34].replace("\n", " ") if has_text
                     else f"<{shp.shape_type}>")

            if x < -0.01 or y < -0.01 or x + w > sw + 0.01 or y + h > sh_in + 0.01:
                problems.append(f"OFF-SLIDE    [{label}] ({x:.2f},{y:.2f}) "
                                f"{w:.2f}x{h:.2f} -> ends ({x + w:.2f},{y + h:.2f})")
            else:
                # A footer band sits closer to the edge by design, and its box is
                # sized for a possible second line so the BOX reaches nearer the
                # edge than the ink does. Judging it by the general 0.4in rule
                # reported five false positives per deck.
                small = has_text and max(
                    (r.font.size.pt for p in shp.text_frame.paragraphs
                     for r in p.runs if r.font.size is not None), default=99) <= 11
                bot_min = FOOTER_MARGIN_IN if (small and y > sh_in * 0.85) else MARGIN_IN
                if (x < MARGIN_IN - 0.01 or y < MARGIN_IN - 0.01
                        or x + w > sw - MARGIN_IN + 0.01
                        or y + h > sh_in - bot_min + 0.01):
                    problems.append(
                        f"TIGHT-MARGIN [{label}] ({x:.2f},{y:.2f}) {w:.2f}x{h:.2f}")

            if shp.shape_type == PICTURE and hasattr(shp, "image"):
                try:
                    px_w, px_h = shp.image.size
                    native, drawn = px_w / px_h, w / h
                    if abs(native - drawn) / native > 0.02:
                        problems.append(
                            f"ASPECT       [image] native {native:.3f} vs drawn {drawn:.3f} "
                            f"({abs(native - drawn) / native * 100:.1f}% distortion)")
                except Exception:
                    pass

            if has_text:
                need = est_text_height(shp, w)
                if need and need > h * 1.08:
                    problems.append(f"TEXT-OVERFLOW[{label}] needs ~{need:.2f}in, "
                                    f"box {h:.2f}in")

            boxes.append((label, box, shp))

        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                (l1, b1, s1), (l2, b2, s2) = boxes[i], boxes[j]
                ov = overlap(b1, b2)
                if ov and s1.has_text_frame and s1.text_frame.text.strip() \
                       and s2.has_text_frame and s2.text_frame.text.strip():
                    problems.append(
                        f"OVERLAP      [{l1}] x [{l2}] by {ov[0]:.2f}x{ov[1]:.2f} in")

        print(f"--- slide {idx:>2} : {len(boxes):>2} shapes : "
              f"{'OK' if not problems else str(len(problems)) + ' ISSUE(S)'}")
        for p in problems:
            print(f"      {p}")
        issues += len(problems)

    print()
    print("=" * 64)
    print(f"TOTAL: {issues} geometric issue(s) across {len(prs.slides)} slides")
    print("=" * 64)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1]))
