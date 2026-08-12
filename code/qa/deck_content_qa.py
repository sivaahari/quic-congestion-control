#!/usr/bin/env python3
"""Content + structure QA for the deck.

Complements deck_geometry_qa.py (which measures boxes). This one re-opens the
package and checks what the XML actually says -- important because the bullet
work writes raw pPr elements that python-pptx has no API for.
"""
import re
import sys
import zipfile

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

# Entity resolution and network access off: stdlib parsers are open to XXE and
# billion-laughs. lxml ships with python-pptx, so this adds no dependency.
SAFE = etree.XMLParser(resolve_entities=False, no_network=True, huge_tree=False)

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def main(path):
    fails = []

    # 1. every part must be well-formed XML and the package must open
    with zipfile.ZipFile(path) as z:
        bad = z.testzip()
        if bad:
            fails.append(f"corrupt zip entry: {bad}")
        for n in z.namelist():
            if n.endswith(".xml") or n.endswith(".rels"):
                try:
                    etree.fromstring(z.read(n), SAFE)
                except etree.XMLSyntaxError as e:
                    fails.append(f"malformed XML in {n}: {e}")
    prs = Presentation(path)
    print(f"opened OK · {len(prs.slides)} slides\n")

    n_bul = n_num = 0
    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shp in slide.shapes:
            if not shp.has_text_frame:
                continue
            for par in shp.text_frame.paragraphs:
                t = "".join(r.text for r in par.runs)
                if not t.strip():
                    continue
                texts.append(t)
                pPr = par._p.find(qn("a:pPr"))
                has_char = pPr is not None and pPr.find(qn("a:buChar")) is not None
                has_num = pPr is not None and pPr.find(qn("a:buAutoNum")) is not None
                n_bul += has_char
                n_num += has_num

                # A literal glyph alongside a real bullet renders twice.
                if has_char and t.lstrip().startswith("•"):
                    fails.append(f"slide {idx}: double bullet -- literal + buChar: {t[:40]!r}")
                if has_num and re.match(r"^\s*\d+[.)]", t):
                    fails.append(f"slide {idx}: double number -- literal + buAutoNum: {t[:40]!r}")
                # Leading whitespace used as indentation does not survive a wrap.
                if t.startswith("  "):
                    fails.append(f"slide {idx}: leading-space indent: {t[:40]!r}")
                # A hanging indent needs marL; without it the wrap falls back.
                if (has_char or has_num) and pPr.get("marL") in (None, "0"):
                    fails.append(f"slide {idx}: bullet without marL: {t[:40]!r}")

        joined = " ".join(texts)
        for pat, why in [(r"\bTODO\b(?! - Correct)", "stray TODO"),
                         (r"\bTBD\b", "TBD"), (r"\bXXX+\b", "placeholder"),
                         (r"lorem|ipsum", "lorem ipsum"),
                         (r"\[missing:", "missing figure"),
                         (r"\bpending\b", "pending marker")]:
            if re.search(pat, joined, re.I):
                fails.append(f"slide {idx}: {why}")
        print(f"  slide {idx:>2}: {len(texts):>2} paragraphs · {len(joined):>4} chars")

    print(f"\nnative bullets: {n_bul} · native numbers: {n_num}")
    print("=" * 60)
    if fails:
        print(f"FAIL — {len(fails)} problem(s)")
        for f in fails:
            print("   " + f)
        return 1
    print("PASS — package well-formed, no double bullets, no placeholders")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1]))
